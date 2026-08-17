import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette - Professional Tech Clean (No Neon)
    COLOR_BG = RGBColor(248, 250, 252)        # Slate 50
    COLOR_CARD = RGBColor(255, 255, 255)      # Pure White
    COLOR_CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    COLOR_PRIMARY = RGBColor(15, 23, 42)      # Slate 900
    COLOR_SECONDARY = RGBColor(71, 85, 105)   # Slate 600
    COLOR_MUTED = RGBColor(100, 116, 139)     # Slate 500
    COLOR_ACCENT = RGBColor(37, 99, 235)      # Blue 600
    COLOR_ACCENT_BG = RGBColor(239, 246, 255) # Blue 50
    COLOR_SUCCESS = RGBColor(22, 163, 74)     # Green 600
    COLOR_SUCCESS_BG = RGBColor(240, 253, 244)# Green 50
    COLOR_CODE_BG = RGBColor(241, 245, 249)   # Slate 100

    def add_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.color.rgb = COLOR_BG
        return bg

    def add_header(slide, title_text, category_text="ITSTEPMARKET • ЕКЗАМЕНАЦІЙНИЙ ПРОЕКТ"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT
        p_cat.font.name = "Segoe UI"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = "Segoe UI"

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==================== SLIDE 1: TITLE SLIDE ====================
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_background(s1)

    card1 = add_card(s1, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5))
    
    tbox = s1.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9.7), Inches(4.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "REACT SPA • ЕКЗАМЕНАЦІЙНА РОБОТА"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.font.name = "Segoe UI"
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "ITSTEPMarket"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Segoe UI"
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Повнофункціональний інтернет-магазин з каталогом, кошиком, збереженням стану та захищеною адмін-панеллю"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_SECONDARY
    p.font.name = "Segoe UI"
    p.space_after = Pt(28)

    p = tf.add_paragraph()
    p.text = "Ключовий стек: React 19 • React Router DOM v7 • Context API + useReducer • LocalStorage • GitHub Pages"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Segoe UI"
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "Автор проекту: Виконано для захисту фінального екзамену IT STEP"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_MUTED
    p.font.name = "Segoe UI"

    # ==================== SLIDE 2: МЕТА ТА ЗАВДАННЯ ====================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_background(s2)
    add_header(s2, "1. Мета, бізнес-логіка та вимоги проекту")

    # 3 Cards
    add_card(s2, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.3))
    t1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.3), Inches(4.9))
    tf1 = t1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Клієнтський модуль"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)
    
    bullets1 = [
        "Каталог товарів з фільтрацією за категоріями та пошуком",
        "Сортування за ціною, рейтингом, назвою",
        "Детальна сторінка товару (/product/:id) з характеристиками",
        "Кошик з можливістю редагування кількості та промокодами",
        "Модальне вікно швидкого оформлення замовлення"
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    add_card(s2, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.3))
    t2 = s2.shapes.add_textbox(Inches(5.0), Inches(1.7), Inches(3.3), Inches(4.9))
    tf2 = t2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Адміністративна панель"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)
    
    bullets2 = [
        "Форма авторизації (/admin/login) з перевіркою пароля",
        "Захист адміністративного маршруту (ProtectedRoute)",
        "Повний CRUD: додавання, перегляд, редагування, видалення",
        "Контрольовані форми з комплексною валідацією полів",
        "Статистичний віджет складу та категорій"
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    add_card(s2, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.3))
    t3 = s2.shapes.add_textbox(Inches(9.0), Inches(1.7), Inches(3.3), Inches(4.9))
    tf3 = t3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Архітектурні вимоги"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)
    
    bullets3 = [
        "Збереження стану кошика та товарів у LocalStorage",
        "Чистий код без зайвих коментарів та сміття",
        "Модульна компонентна структура (Separation of Concerns)",
        "Стриманий дизайн без неонових та перевантажених елементів",
        "Повна сумісність з GitHub Pages (HashRouter + deploy)"
    ]
    for b in bullets3:
        p = tf3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 3: ТЕХНОЛОГІЧНИЙ СТЕК ====================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_background(s3)
    add_header(s3, "2. Технологічний стек та обґрунтування вибору")

    stack_items = [
        ("React 19", "Сучасна основа SPA", "Використання хуків useState, useReducer, useContext, useMemo, useEffect для побудови реактивного інтерфейсу без надлишкових зовнішніх бібліотек."),
        ("React Router DOM v7", "Клієнтська маршрутизація", "Організація навігації між сторінками. Використано HashRouter для 100% сумісності та виключення помилок 404 на GitHub Pages."),
        ("Vite 8", "Блискавичний бандлер", "Миттєвий HMR під час розробки та високооптимізована продакшн-збірка за 470 мс з мініфікацією та tree-shaking."),
        ("Native Context + useReducer", "State Management", "Централізоване управління даними магазину (товари, кошик, авторизація) за патерном Redux без сторонніх залежностей."),
        ("Vanilla CSS & Tokens", "Дизайн-система", "CSS Custom Properties для кольорів та відступів, гнучкі Flexbox/Grid лейаути, легкість у підтримці, чиста естетика без зайвого неону."),
        ("LocalStorage API", "Шар персистентності", "Синхронізація клієнтських мутацій з браузерним сховищем для збереження товарів та стану кошика між перезавантаженнями.")
    ]

    for idx, (title, subtitle, desc) in enumerate(stack_items):
        row = idx // 3
        col = idx % 3
        left = Inches(0.8 + col * 4.0)
        top = Inches(1.5 + row * 2.7)
        add_card(s3, left, top, Inches(3.7), Inches(2.4))
        
        tb = s3.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY

    # ==================== SLIDE 4: ЗАГАЛЬНА АРХІТЕКТУРА ДЕРЕВА КОМПОНЕНТІВ ====================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_background(s4)
    add_header(s4, "3. Архітектура компонентів та ієрархія провайдерів")

    # Left: Providers tree
    add_card(s4, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_tree = s4.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_tree = tb_tree.text_frame
    tf_tree.word_wrap = True
    
    p = tf_tree.paragraphs[0]
    p.text = "Ієрархія провайдерів (App.jsx)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    tree_lines = [
        "<AuthProvider> // Стан авторизації адміністратора",
        "  └── <ProductProvider> // Каталог та CRUD товарів",
        "        └── <CartProvider> // Кошик, знижки, розрахунки",
        "              └── <HashRouter> // Маршрутизація SPA",
        "                    └── <Layout> // Header, Footer, Toast",
        "                          └── <Routes>",
        "                                ├── / (CatalogPage)",
        "                                ├── /product/:id (ProductDetailPage)",
        "                                ├── /cart (CartPage)",
        "                                ├── /admin/login (AdminLoginPage)",
        "                                ├── /admin (<ProtectedRoute><AdminDashboardPage/></ProtectedRoute>)",
        "                                └── * (NotFoundPage)"
    ]
    for tl in tree_lines:
        p = tf_tree.add_paragraph()
        p.text = tl
        p.font.size = Pt(10.5)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY if not tl.strip().startswith("//") else COLOR_MUTED
        p.space_after = Pt(2)

    # Right: Description
    add_card(s4, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_desc = s4.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_desc = tb_desc.text_frame
    tf_desc.word_wrap = True
    
    p = tf_desc.paragraphs[0]
    p.text = "Ключові архітектурні рішення"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    arch_bullets = [
        ("Вкладеність провайдерів (Provider Tree)", "Провайдери розташовані відповідно до залежностей: Auth не залежить від кошика, Product надає товари для Cart, а Cart надає дані у глобальний інтерфейс."),
        ("Глобальний Layout", "Компонент Layout загортає всі сторінки, забезпечуючи стабільну шапку з лічильником кошика, футер та глобальну шину сповіщень (NotificationToast)."),
        ("Ізоляція маршрутів", "Кожна сторінка є автономним модулем, що отримує дані виключно через власні кастомні хуки (useAuth, useProducts, useCart)."),
        ("Single Responsibility Principle", "Чіткий поділ на UI-компоненти (чисті відображення) та Context/State logic (бізнес-правила та персистентність).")
    ]
    for title, text in arch_bullets:
        p = tf_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 5: СТРУКТУРА ПРОЕКТУ ====================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_background(s5)
    add_header(s5, "4. Структура файлової системи (Clean Architecture)")

    folders = [
        ("src/context/", "Керування глобальним станом", "• AuthContext.jsx — сесія адміністратора\n• ProductContext.jsx — каталог товарів (CRUD)\n• CartContext.jsx — операції над кошиком та сповіщення"),
        ("src/pages/", "Сторінки маршрутизатора", "• CatalogPage — головна вітрина з пошуком\n• ProductDetailPage — деталі обраного товару\n• CartPage — кошик та оформлення замовлення\n• AdminLoginPage / AdminDashboardPage — адмінка\n• NotFoundPage — сторінка 404"),
        ("src/components/", "Модульні UI компоненти", "• /admin — AdminProductTable, ProductFormModal\n• /cart — CartItem, CartSummary\n• /product — ProductCard, ProductGrid, ProductFilter\n• /layout — Header, Footer, Layout\n• /common — ProtectedRoute, NotificationToast"),
        ("src/styles/ & data/", "Дизайн-система та утиліти", "• variables.css — дизайн-токени (кольори, радіуси)\n• global.css — скидання стилів та базові класи\n• mockProducts.js — початкова база товарів\n• formatters.js — форматування цін (гривні)")
    ]

    for idx, (title, sub, details) in enumerate(folders):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.5 + row * 2.7)
        add_card(s5, left, top, Inches(5.6), Inches(2.5))
        
        tb = s5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(4)
        
        for line in details.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_SECONDARY

    # ==================== SLIDE 6: УПРАВЛІННЯ СТАНОМ ====================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_background(s6)
    add_header(s6, "5. Управління глобальним станом: Context API + useReducer")

    add_card(s6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_code = s6.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = True
    
    p = tf_code.paragraphs[0]
    p.text = "Приклад реалізації CartReducer"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    reducer_code = [
        "const cartReducer = (state, action) => {",
        "  switch (action.type) {",
        "    case 'ADD_TO_CART': {",
        "      const { product, quantity } = action.payload;",
        "      const idx = state.findIndex(",
        "        item => item.product.id === product.id",
        "      );",
        "      if (idx > -1) {",
        "        return state.map((item, i) => ",
        "          i === idx ? { ...item, quantity: item.quantity + quantity } : item",
        "        );",
        "      }",
        "      return [...state, { product, quantity }];",
        "    }",
        "    case 'UPDATE_QUANTITY':",
        "      return state.map(item =>",
        "        item.product.id === action.payload.productId",
        "          ? { ...item, quantity: action.payload.quantity }",
        "          : item",
        "      );",
        "    case 'REMOVE_FROM_CART':",
        "      return state.filter(i => i.product.id !== action.payload);",
        "  }",
        "};"
    ]
    for line in reducer_code:
        p = tf_code.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY

    add_card(s6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_why = s6.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_why = tb_why.text_frame
    tf_why.word_wrap = True
    
    p = tf_why.paragraphs[0]
    p.text = "Чому useReducer + Context, а не Redux?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    reasons = [
        ("Передбачуваність мутацій (Pure Reducer)", "Стан не мутується напряму. Будь-які зміни відбуваються через чисті функції без сайд-ефектів, що гарантує надійність додатка."),
        ("Інкапсуляція дій в кастомні хуки", "Компоненти викликають зрозумілі методи: addToCart(prod), updateQuantity(id, qty), deleteProduct(id) без необхідності знати внутрішній формат dispatch."),
        ("Нульова вага додаткових бібліотек", "Використовуються нативні інструменти React 19 без стороннього бойлерплейту (Redux Toolkit, Zustand), що зменшує розмір фінального бандла."),
        ("Атомарні контексти", "Розподіл на 3 незалежні контексти (Auth, Products, Cart) виключає зайві ре-рендери компонентів при зміні непов'язаних даних.")
    ]
    for title, text in reasons:
        p = tf_why.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_why.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 7: КАТАЛОГ ТА ФІЛЬТРАЦІЯ ====================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_background(s7)
    add_header(s7, "6. Каталог товарів: Алгоритми пошуку, фільтрації та сортування")

    add_card(s7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_c1 = s7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "Оптимізована фільтрація через useMemo"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    filter_code = [
        "const filteredAndSortedProducts = useMemo(() => {",
        "  return products",
        "    .filter(product => {",
        "      const matchesSearch =",
        "        product.title.toLowerCase().includes(searchTerm.toLowerCase().trim()) ||",
        "        product.description.toLowerCase().includes(searchTerm.toLowerCase().trim());",
        "      const matchesCategory =",
        "        selectedCategory === 'Усі категорії' ||",
        "        product.category === selectedCategory;",
        "      return matchesSearch && matchesCategory;",
        "    })",
        "    .sort((a, b) => {",
        "      if (sortBy === 'price-asc') return a.price - b.price;",
        "      if (sortBy === 'price-desc') return b.price - a.price;",
        "      if (sortBy === 'rating') return (b.rating || 0) - (a.rating || 0);",
        "      if (sortBy === 'name') return a.title.localeCompare(b.title, 'uk');",
        "      return 0;",
        "    });",
        "}, [products, searchTerm, selectedCategory, sortBy]);"
    ]
    for line in filter_code:
        p = tf_c1.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY

    add_card(s7, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_c2 = s7.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    
    p = tf_c2.paragraphs[0]
    p.text = "Особливості реалізації каталогу"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    cat_feats = [
        ("Двосторонній комбінований пошук", "Пошуковий запит одночасно шукає збіги за заголовком товару та за його детальним описом з авто-очищенням пробілів (trim)."),
        ("Локалізоване сортування (localeCompare)", "Сортування за алфавітом використовує українську локаль 'uk', що забезпечує коректний порядок літер (А, Б, В, Г, Ґ...)."),
        ("Мемоізація обчислень", "Масив фільтрується виключно при зміні критеріїв фільтрації або списку товарів завдяки useMemo."),
        ("Стан кнопки 'В кошику'", "Картка автоматично визначає, чи є вже цей товар у кошику (isInCart), та змінює текст і стан кнопки покупки.")
    ]
    for title, text in cat_feats:
        p = tf_c2.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_c2.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 8: СТОРІНКА ТОВАРУ ====================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_background(s8)
    add_header(s8, "7. Детальна сторінка товару (/product/:id)")

    add_card(s8, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_p1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_p1 = tb_p1.text_frame
    tf_p1.word_wrap = True
    
    p = tf_p1.paragraphs[0]
    p.text = "Отримання параметрів та захист"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    prod_code = [
        "export const ProductDetailPage = () => {",
        "  const { id } = useParams();",
        "  const navigate = useNavigate();",
        "  const { getProductById } = useProducts();",
        "  const { addToCart, cartItems } = useCart();",
        "  const [quantity, setQuantity] = useState(1);",
        "",
        "  const product = getProductById(id);",
        "",
        "  if (!product) {",
        "    return (",
        "      <div className=\"product-not-found\">",
        "        <h2>Товар не знайдено</h2>",
        "        <Link to=\"/\">Повернутися</Link>",
        "      </div>",
        "    );",
        "  }",
        "  // Рендеринг галереї, цін, характеристик",
        "};"
    ]
    for line in prod_code:
        p = tf_p1.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY

    add_card(s8, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_p2 = s8.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True
    
    p = tf_p2.paragraphs[0]
    p.text = "Функціонал картки товару"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    prod_feats = [
        ("Динамічний роутинг useParams", "Вилучення ідентифікатора товару з URL та пошук відповідного об'єкта через Context без повторних мережевих запитів."),
        ("Обробка неіснуючого товару (Fallback)", "Якщо id в URL некоректний або товар був видалений адміном, користувач бачить інформативне повідомлення з кнопкою повернення."),
        ("Селектор кількості (Quantity Controls)", "Можливість обрати потрібну кількість одиниць товару перед додаванням у кошик із захистом від від'ємних значень (Math.max(1, q-1))."),
        ("Динамічні технічні характеристики", "Гнучкий рендеринг об'єкта specs у вигляді адаптивної таблиці характеристик (екран, процесор, пам'ять, акумулятор тощо).")
    ]
    for title, text in prod_feats:
        p = tf_p2.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_p2.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 9: СИСТЕМА КОШИКА ====================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_background(s9)
    add_header(s9, "8. Модуль кошика: Операції, лічильники та інтеграція")

    add_card(s9, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_cart_ops = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_cart_ops = tb_cart_ops.text_frame
    tf_cart_ops.word_wrap = True
    
    p = tf_cart_ops.paragraphs[0]
    p.text = "Основні операції над елементами кошика"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    cart_operations = [
        ("Додавання товару (addToCart)", "При першому додаванні створюється новий елемент { product, quantity }. При повторному — збільшується існуюче поле quantity без дублювання рядків."),
        ("Зміна кількості (updateQuantity)", "Кнопки '+' та '-' дозволяють змінювати кількість безпосередньо в таблиці кошика з блокуванням кнопки '-' при кількості 1."),
        ("Видалення товару (removeFromCart)", "Повне видалення позиції за product.id з відображенням спливаючого сповіщення через Toast Notification."),
        ("Очищення кошика (clearCart)", "Скидання масиву товарів до [] після успішного оформлення замовлення або за бажанням покупця.")
    ]
    for title, text in cart_operations:
        p = tf_cart_ops.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_cart_ops.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    add_card(s9, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_cart_code = s9.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_cart_code = tb_cart_code.text_frame
    tf_cart_code.word_wrap = True
    
    p = tf_cart_code.paragraphs[0]
    p.text = "Агрегація даних у CartContext"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    calc_code = [
        "// Автоматичний підрахунок вартості",
        "const totalAmount = cartItems.reduce(",
        "  (sum, item) => sum + item.product.price * item.quantity,",
        "  0",
        ");",
        "",
        "// Загальна кількість одиниць (для бейджа в Header)",
        "const totalItemsCount = cartItems.reduce(",
        "  (sum, item) => sum + item.quantity,",
        "  0",
        ");",
        "",
        "// Toast сповіщення",
        "const showToast = (message) => {",
        "  setToastMessage(message);",
        "  setTimeout(() => setToastMessage(null), 3000);",
        "};"
    ]
    for line in calc_code:
        p = tf_cart_code.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY if not line.startswith("//") else COLOR_MUTED

    # ==================== SLIDE 10: ОФОРМЛЕННЯ ЗАМОВЛЕННЯ ====================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_background(s10)
    add_header(s10, "9. Розрахунок замовлення, промокоди та Checkout Flow")

    add_card(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_sum1 = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_sum1 = tb_sum1.text_frame
    tf_sum1.word_wrap = True
    
    p = tf_sum1.paragraphs[0]
    p.text = "Система знижок та промокодів"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    promo_points = [
        ("Механіка промокоду", "У компоненті CartSummary реалізовано обробку промокоду 'TECH10', який надає миттєву знижку 10% на все замовлення."),
        ("Індикація стану промокоду", "Відображення статусу застосування (зелений індикатор успіху або червоне попередження про недійсний код)."),
        ("Форматування валюти (formatPrice)", "Єдина утиліта formatPrice(val) форматує числа з розділювачами тисяч у стандартному українському форматі (наприклад: 48 999 ₴).")
    ]
    for title, text in promo_points:
        p = tf_sum1.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_sum1.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    add_card(s10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_sum2 = s10.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_sum2 = tb_sum2.text_frame
    tf_sum2.word_wrap = True
    
    p = tf_sum2.paragraphs[0]
    p.text = "Checkout Модальне вікно"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    checkout_points = [
        ("Контрольовані поля форми", "Збір даних покупця: ПІБ, валідований телефон (перевірка формату regex), адреса відділення та вибір способу оплати."),
        ("Генерація номера замовлення", "При успішному сабміті генерується унікальний 6-значний номер замовлення (#XXXXXX) та відображається підсумок."),
        ("Автоматичне очищення", "Після оформлення замовлення кошик автоматично очищується, а клієнт отримує інструкцію щодо зв'язку з менеджером.")
    ]
    for title, text in checkout_points:
        p = tf_sum2.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_sum2.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 11: ЗБЕРЕЖЕННЯ ДАНИХ (LOCALSTORAGE) ====================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_background(s11)
    add_header(s11, "10. Шар збереження даних (LocalStorage Persistence)")

    add_card(s11, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_ls_code = s11.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_ls_code = tb_ls_code.text_frame
    tf_ls_code.word_wrap = True
    
    p = tf_ls_code.paragraphs[0]
    p.text = "Двоетапна персистентність"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    ls_code = [
        "// 1. Лінива ініціалізація стану з LocalStorage",
        "const [products, dispatch] = useReducer(",
        "  productReducer,",
        "  [],",
        "  () => {",
        "    const saved = localStorage.getItem('shop_products');",
        "    if (saved) {",
        "      try { return JSON.parse(saved); }",
        "      catch (err) { console.error(err); }",
        "    }",
        "    return INITIAL_PRODUCTS; // Fallback",
        "  }",
        ");",
        "",
        "// 2. Реактивна синхронізація при будь-яких змінах",
        "useEffect(() => {",
        "  localStorage.setItem('shop_products', JSON.stringify(products));",
        "}, [products]);"
    ]
    for line in ls_code:
        p = tf_ls_code.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY if not line.startswith("//") else COLOR_MUTED

    add_card(s11, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_ls_desc = s11.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_ls_desc = tb_ls_desc.text_frame
    tf_ls_desc.word_wrap = True
    
    p = tf_ls_desc.paragraphs[0]
    p.text = "Чому це надійне технічне рішення?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    ls_advantages = [
        ("Lazy Initializer у useReducer", "Зчитування з localStorage виконується рівно один раз під час монтажу компонента, виключаючи блокування головного потоку при повторних рендерах."),
        ("Захист через try-catch", "Якщо дані в сховищі пошкоджені або заблоковані політикою браузера, додаток автоматично завантажує мокові дані INITIAL_PRODUCTS без падіння."),
        ("Синхронізація товарів і кошика", "Як клієнтський кошик ('shop_cart'), так і товари ('shop_products') та сесія адміна ('shop_admin_auth') зберігаються автономно."),
        ("Повна автономність", "Додаток працює як повноцінний автономний SPA без потреби підняття важкого бекенду для демонстрації.")
    ]
    for title, text in ls_advantages:
        p = tf_ls_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_ls_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 12: АВТОРИЗАЦІЯ АДМІНІСТРАТОРА ====================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_background(s12)
    add_header(s12, "11. Модуль авторизації адміністратора (AuthContext)")

    add_card(s12, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_auth_code = s12.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_auth_code = tb_auth_code.text_frame
    tf_auth_code.word_wrap = True
    
    p = tf_auth_code.paragraphs[0]
    p.text = "Логіка авторизації (AuthContext.jsx)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    auth_code = [
        "export const AuthProvider = ({ children }) => {",
        "  const [isAdminLoggedIn, setIsAdminLoggedIn] = useState(() => {",
        "    return localStorage.getItem('shop_admin_auth') === 'true';",
        "  });",
        "",
        "  const login = (username, password) => {",
        "    if (username === 'admin' && password === 'admin123') {",
        "      setIsAdminLoggedIn(true);",
        "      return { success: true };",
        "    }",
        "    return {",
        "      success: false,",
        "      error: 'Невірний логін або пароль.'",
        "    };",
        "  };",
        "",
        "  const logout = () => setIsAdminLoggedIn(false);",
        "  // Повернення провайдера",
        "};"
    ]
    for line in auth_code:
        p = tf_auth_code.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY

    add_card(s12, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_auth_desc = s12.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_auth_desc = tb_auth_desc.text_frame
    tf_auth_desc.word_wrap = True
    
    p = tf_auth_desc.paragraphs[0]
    p.text = "Особливості реалізації безпеки"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    auth_feats = [
        ("Облікові дані адміністратора", "Логін: admin | Пароль: admin123. Передбачено захист від випадкових пробілів через String.trim()."),
        ("Синхронізація сесії", "Стан авторизації зберігається в localStorage, тому після оновлення сторінки сесія адміністратора не скидається."),
        ("Кнопка швидкого виходу (Logout)", "У шапці сайту (Header) для авторизованого адміна відображається кнопка переходу до панелі та кнопка виходу з системи."),
        ("Індикація помилок", "При введенні некоректних даних виводиться стилізований алерт із підказкою.")
    ]
    for title, text in auth_feats:
        p = tf_auth_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_auth_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 13: ЗАХИСТ МАРШРУТІВ (PROTECTED ROUTE) ====================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_background(s13)
    add_header(s13, "12. Захист маршрутів (ProtectedRoute Guard Pattern)")

    add_card(s13, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_guard = s13.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_guard = tb_guard.text_frame
    tf_guard.word_wrap = True
    
    p = tf_guard.paragraphs[0]
    p.text = "Компонент ProtectedRoute"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    guard_code = [
        "import { Navigate } from 'react-router-dom';",
        "import { useAuth } from '../../context/AuthContext';",
        "",
        "export const ProtectedRoute = ({ children }) => {",
        "  const { isAdminLoggedIn } = useAuth();",
        "",
        "  if (!isAdminLoggedIn) {",
        "    // Автоматичний редирект на вхід",
        "    return <Navigate to=\"/admin/login\" replace />;",
        "  }",
        "",
        "  // Доступ дозволено",
        "  return children;",
        "};"
    ]
    for line in guard_code:
        p = tf_guard.add_paragraph()
        p.text = line
        p.font.size = Pt(10.5)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY if not line.startswith("//") else COLOR_MUTED

    add_card(s13, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_guard_desc = s13.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_guard_desc = tb_guard_desc.text_frame
    tf_guard_desc.word_wrap = True
    
    p = tf_guard_desc.paragraphs[0]
    p.text = "Переваги шаблону ProtectedRoute"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    guard_advantages = [
        ("Декларативний захист у роутері", "Маршрут загортається у <ProtectedRoute><AdminDashboardPage/></ProtectedRoute> безпосередньо в конфігурації App.jsx."),
        ("Заміна історії переходів (replace)", "Параметр replace у компоненті Navigate запобігає зацикленню переходів при натисканні кнопки 'Назад' у браузері."),
        ("Миттєве блокування неавторизованих запитів", "Користувач без прав не може побачити інтерфейс управління товарами, змінити залишки чи видалити позиції."),
        ("Зворотний редирект для авторизованих", "Якщо вже залогінений адмін переходить на /admin/login, сторінка автоматично перенаправляє його на /admin.")
    ]
    for title, text in guard_advantages:
        p = tf_guard_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_guard_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 14: УПРАВЛІННЯ ТОВАРАМИ (CRUD) ====================
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_background(s14)
    add_header(s14, "13. Адміністративна панель: Управління товарами (CRUD)")

    crud_cards = [
        ("CREATE (Створення)", "Додавання нового товару через модальне вікно ProductFormModal. Автоматична генерація унікального id (Date.now()), дефолтного рейтингу та додавання на початок списку."),
        ("READ (Відображення)", "Таблиця AdminProductTable з мініатюрами фото, статусом наявності та кількістю на складі. Швидкий пошук та статистика (загальна кількість, категорії, залишки)."),
        ("UPDATE (Редагування)", "Підвантаження існуючих даних товару у контрольовану форму з можливістю змінити ціну, назву, категорію, кількість на складі, опис та бейджі ('Новинка', 'Хіт')."),
        ("DELETE (Видалення)", "Безпечне видалення з підтвердженням через спеціальне діалогове вікно попередження, що унеможливлює випадкове стирання позицій.")
    ]

    for idx, (title, desc) in enumerate(crud_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.5 + row * 2.7)
        add_card(s14, left, top, Inches(5.6), Inches(2.5))
        
        tb = s14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SECONDARY

    # ==================== SLIDE 15: АРХІТЕКТУРА ФОРМ ТА ВАЛІДАЦІЯ ====================
    s15 = prs.slides.add_slide(blank_layout)
    add_slide_background(s15)
    add_header(s15, "14. Контрольовані форми та комплексна валідація")

    add_card(s15, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_val_code = s15.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_val_code = tb_val_code.text_frame
    tf_val_code.word_wrap = True
    
    p = tf_val_code.paragraphs[0]
    p.text = "Валідація полів (ProductFormModal.jsx)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    val_code = [
        "const validate = () => {",
        "  const newErrors = {};",
        "  if (!formData.title.trim()) {",
        "    newErrors.title = 'Введіть назву товару';",
        "  } else if (formData.title.trim().length < 3) {",
        "    newErrors.title = 'Мінімум 3 символи';",
        "  }",
        "  if (!formData.price || Number(formData.price) <= 0) {",
        "    newErrors.price = 'Ціна повинна бути додатним числом';",
        "  }",
        "  if (!formData.image.startsWith('http://') &&",
        "      !formData.image.startsWith('https://')) {",
        "    newErrors.image = 'URL має починатися з http:// або https://';",
        "  }",
        "  if (formData.description.trim().length < 10) {",
        "    newErrors.description = 'Опис не менше 10 символів';",
        "  }",
        "  setErrors(newErrors);",
        "  return Object.keys(newErrors).length === 0;",
        "};"
    ]
    for line in val_code:
        p = tf_val_code.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY

    add_card(s15, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_val_desc = s15.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_val_desc = tb_val_desc.text_frame
    tf_val_desc.word_wrap = True
    
    p = tf_val_desc.paragraphs[0]
    p.text = "Принципи роботи з формами"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    val_rules = [
        ("Controlled Components", "Кожне поле (input, select, textarea, checkbox) прив'язане до стану formData через єдиний обробник handleChange."),
        ("Миттєве скидання помилки при вводі", "Коли користувач починає виправляти поле, помилка для цього поля автоматично видаляється зі стану errors."),
        ("Запобігання відправки невалідних даних", "Функція handleSubmit блокує виклик onSave, якщо валідація не пройшла."),
        ("Типізація та санітизація", "Ціни та залишки приводяться до типу Number, а текстові рядки очищуються методом trim().")
    ]
    for title, text in val_rules:
        p = tf_val_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_val_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 16: CSS АРХІТЕКТУРА ТА ДИЗАЙН ====================
    s16 = prs.slides.add_slide(blank_layout)
    add_slide_background(s16)
    add_header(s16, "15. CSS Архітектура, адаптивність та дизайн-система")

    css_cards = [
        ("Дизайн-токени (variables.css)", "Єдині змінні для кольорів (--primary-color, --bg-card, --border-color), радіусів (--radius-md: 6px) та плавності анімацій (--transition-fast: 0.1s)."),
        ("Чиста естетика (Без неону)", "Стриманий сучасний мінімалістичний стиль. Чітка типографіка (Inter/Segoe UI), збалансовані відступи та висока контрастність для комфорту очей."),
        ("Адаптивність (Mobile First & Breakpoints)", "Адаптивні медіа-запити (@media max-width: 768px, 1024px) забезпечують бездоганне відображення на смартфонах, планшетах і ПК."),
        ("Модульність стилів", "Кожен компонент та сторінка має власний CSS файл (ProductCard.css, CartPage.css тощо), що спрощує підтримку та усуває конфлікти селекторів.")
    ]

    for idx, (title, desc) in enumerate(css_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.5 + row * 2.7)
        add_card(s16, left, top, Inches(5.6), Inches(2.5))
        
        tb = s16.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_SECONDARY

    # ==================== SLIDE 17: ЗБІРКА ТА ДЕПЛОЙ НА GITHUB PAGES ====================
    s17 = prs.slides.add_slide(blank_layout)
    add_slide_background(s17)
    add_header(s17, "16. Продакшн-збірка та автоматизований деплой на GitHub Pages")

    add_card(s17, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_dep_code = s17.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_dep_code = tb_dep_code.text_frame
    tf_dep_code.word_wrap = True
    
    p = tf_dep_code.paragraphs[0]
    p.text = "Конфігурація Vite та package.json"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(8)

    dep_code = [
        "// vite.config.js — Відносні шляхи для GitHub Pages",
        "export default defineConfig({",
        "  plugins: [react()],",
        "  base: './'",
        "});",
        "",
        "// package.json — Скрипти деплою",
        "\"scripts\": {",
        "  \"dev\": \"vite\",",
        "  \"build\": \"vite build\",",
        "  \"predeploy\": \"npm run build\",",
        "  \"deploy\": \"gh-pages -d dist\"",
        "}"
    ]
    for line in dep_code:
        p = tf_dep_code.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_PRIMARY if not line.startswith("//") else COLOR_MUTED

    add_card(s17, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_dep_desc = s17.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_dep_desc = tb_dep_desc.text_frame
    tf_dep_desc.word_wrap = True
    
    p = tf_dep_desc.paragraphs[0]
    p.text = "Кроки та вирішення проблем деплою"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(10)

    deploy_steps = [
        ("Чому base: './'?", "Забезпечує коректне підвантаження зібраних CSS та JS файлів незалежно від назви репозиторію на github.io."),
        ("Чому HashRouter?", "На статичному хостингу GitHub Pages сервер не підтримує динамічний rewrite маршрутів. HashRouter повністю вирішує проблему 404 помилок при оновленні сторінок."),
        ("Однокомандний деплой", "Команда 'npm run deploy' автоматично викликає predeploy (збірку проекту у dist) та публікує результат у гілку gh-pages."),
        ("Мінімальний розмір збірки", "Завдяки оптимізації Vite розмір стилів складає всього 19.5 kB, а JS-бандла — 88 kB (gzip).")
    ]
    for title, text in deploy_steps:
        p = tf_dep_desc.add_paragraph()
        p.text = "• " + title + ":"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p = tf_dep_desc.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # ==================== SLIDE 18: ПИТАННЯ ТА ВІДПОВІДІ ДЛЯ ЕКЗАМЕНАТОРА ====================
    s18 = prs.slides.add_slide(blank_layout)
    add_slide_background(s18)
    add_header(s18, "17. Захист проекту: Відповіді на типові питання екзаменатора (Q&A)")

    add_card(s18, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3))
    tb_qa1 = s18.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_qa1 = tb_qa1.text_frame
    tf_qa1.word_wrap = True
    
    qa_list1 = [
        ("П: Чому не використано бібліотеку валідації (наприклад, Yup / Zod)?",
         "В: Для контрольованих форм магазину нативна валідація через JS регулярні вирази є легковагою, прозорою, не обтяжує бандл та демонструє глибоке розуміння роботи зі станом у React."),
        ("П: Як уникнути Prop Drilling у великих додатках?",
         "В: У проекті цю проблему повністю вирішено за допомогою Context API. Компоненти отримують доступ до даних напряму через кастомні хуки (useCart, useProducts, useAuth) на будь-якому рівні вкладеності.")
    ]
    for q, a in qa_list1:
        p = tf_qa1.add_paragraph()
        p.text = q
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(4)
        p = tf_qa1.add_paragraph()
        p.text = a
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(14)

    add_card(s18, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3))
    tb_qa2 = s18.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_qa2 = tb_qa2.text_frame
    tf_qa2.word_wrap = True
    
    qa_list2 = [
        ("П: Як підключити реальний REST API бекенд?",
         "В: Архітектура Context готова до заміни LocalStorage на асинхронні запити (fetch/axios). Достатньо всередині addProduct, editProduct, deleteProduct надсилати HTTP запити (POST, PUT, DELETE) до API та оновлювати стан."),
        ("П: Як забезпечено продуктивність інтерфейсу?",
         "В: Застосовано useMemo для фільтрації списку товарів, ліниву ініціалізацію стану useReducer з LocalStorage, оптимізацію рендеру списків через унікальні key та легковагий Lucide React для іконок.")
    ]
    for q, a in qa_list2:
        p = tf_qa2.add_paragraph()
        p.text = q
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(4)
        p = tf_qa2.add_paragraph()
        p.text = a
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(14)

    output_path = "d:/ITSTEP/React/finalproject/ITSTEPMarket_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
